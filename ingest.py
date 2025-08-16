# ingest.py
from __future__ import annotations
from pathlib import Path
from typing import Tuple, List, Callable
import tempfile

import docx2txt
from unstructured.partition.auto import partition
from image_utils import resize_images_batch
import fitz
from pptx import Presentation
from charset_normalizer import from_bytes

SUPPORTED_EXTS = {
    ".jpg",
    ".jpeg",
    ".png",
    # ".doc",
    ".docx",
    ".odt",
    ".pdf",
    ".ppt",
    ".pptx",
    ".txt",
}
MAX_TEXT_CHARS = 20000
MAX_IMAGES = 3
MAX_IMAGE_DIMENSION = 1024  # Maximale Bilddimension in Pixeln
PDF_RASTER_DPI = 200


def is_supported_file(p: Path) -> bool:
    return p.is_file() and p.suffix.lower() in SUPPORTED_EXTS


def _truncate(text: str, limit: int = MAX_TEXT_CHARS) -> str:
    text = (text or "").strip()
    return text if len(text) <= limit else text[:limit] + "…"


def _read_text_safely(path: Path) -> str:
    """
    Liest Bytes und versucht:
      1) UTF-8-SIG (BOM wird entfernt),
      2) charset-normalizer (best guess),
      3) latin-1 (verlustfrei, nie DecodeError).
    """
    data = path.read_bytes()
    # 1) Bevorzugt UTF-8 (mit/ohne BOM)
    try:
        return data.decode("utf-8-sig")
    except UnicodeDecodeError:
        pass

    # 2) Automatische Erkennung
    try:
        match = from_bytes(data).best()
        if match and match.encoding:
            return data.decode(match.encoding, errors="replace")
    except Exception:
        pass

    # 3) Fallback: latin-1 (verlustfrei)
    return data.decode("latin-1", errors="replace")


def ingest_image(path: Path) -> Tuple[str, List[str], callable]:
    # Bild auf maximal MAX_IMAGE_DIMENSION skalieren und an das Modell geben
    resized_paths, cleanup = resize_images_batch(
        [str(path)], max_dimension=MAX_IMAGE_DIMENSION
    )
    return "", resized_paths, cleanup


def ingest_txt(path: Path) -> Tuple[str, List[str], Callable]:
    text = _read_text_safely(path)
    text = text.replace("\r\n", "\n").replace("\r", "\n").strip()
    return _truncate(text), [], (lambda: None)


def ingest_docx(path: Path) -> Tuple[str, List[str], callable]:
    # Bilder via docx2txt in temp-Verzeichnis extrahieren
    tmpdir_obj = tempfile.TemporaryDirectory(prefix="docsort_docx_")
    tmpdir = tmpdir_obj.name
    text = docx2txt.process(str(path), tmpdir)
    images = sorted(Path(tmpdir).glob("*"))
    img_paths = [str(p) for p in images[:MAX_IMAGES]]

    # Extrahierte Bilder skalieren
    resized_paths, resize_cleanup = resize_images_batch(
        img_paths, max_dimension=MAX_IMAGE_DIMENSION
    )

    # Kombinierte Cleanup-Funktion
    def combined_cleanup():
        try:
            resize_cleanup()
        except Exception:
            pass
        try:
            tmpdir_obj.cleanup()
        except Exception:
            pass

    return _truncate(text), resized_paths, combined_cleanup


def ingest_via_unstructured(path: Path) -> Tuple[str, List[str], callable]:
    els = partition(filename=str(path))
    text = _truncate("\n".join(e.text for e in els if getattr(e, "text", None)))
    return text, [], (lambda: None)


def ingest_pdf(path: Path) -> Tuple[str, List[str], Callable]:
    """
    Extrahiert Text mit PyMuPDF und rastert die ersten N Seiten (MAX_IMAGES) in PNGs.
    """
    tmpdir_obj = tempfile.TemporaryDirectory(prefix="docsort_pdf_")
    tmpdir = Path(tmpdir_obj.name)
    texts: List[str] = []
    images: List[str] = []

    try:
        doc = fitz.open(str(path))
        for idx, page in enumerate(doc):
            # Text sammeln
            try:
                texts.append(page.get_text("text"))  # nativer Text-Extractor
            except Exception:
                pass
            # Erste N Seiten als PNG rendern
            if idx < MAX_IMAGES:
                pix = page.get_pixmap(
                    dpi=PDF_RASTER_DPI
                )  # schnelles, direktes Rasterisieren
                out = tmpdir / f"page_{idx+1}.png"
                pix.save(str(out))
                images.append(str(out))
            # Stop-Heuristik gegen sehr lange PDFs
            if sum(len(t) for t in texts) > MAX_TEXT_CHARS * 1.5:
                break
    except Exception:
        # Fallback: unstructured
        return ingest_via_unstructured(path)

    # Extrahierte PNG-Seiten skalieren
    resized_paths, resize_cleanup = resize_images_batch(
        images, max_dimension=MAX_IMAGE_DIMENSION
    )

    # Kombinierte Cleanup-Funktion
    def combined_cleanup():
        try:
            resize_cleanup()
        except Exception:
            pass
        try:
            tmpdir_obj.cleanup()
        except Exception:
            pass

    excerpt = _truncate("\n".join(texts))
    return excerpt, resized_paths, combined_cleanup


def ingest_pptx(path: Path) -> Tuple[str, List[str], Callable]:
    """
    Extrahiert pro Folie Titel, Notizen und Bilder (falls vorhanden).
    """
    try:
        prs = Presentation(str(path))
    except Exception:
        # Fallback
        return ingest_via_unstructured(path)

    pieces: List[str] = []
    images: List[str] = []
    tmpdir_obj = tempfile.TemporaryDirectory(prefix="docsort_pptx_")
    tmpdir = Path(tmpdir_obj.name)

    try:
        for i, slide in enumerate(prs.slides, start=1):
            # Titel
            title = None
            if slide.shapes.title is not None:
                try:
                    title = slide.shapes.title.text
                except Exception:
                    title = None

            # Body-Text (optional, knapp)
            body_chunks: List[str] = []
            for shp in slide.shapes:
                if getattr(shp, "has_text_frame", False):
                    try:
                        txt = (shp.text or "").strip()
                        if txt:
                            body_chunks.append(txt)
                    except Exception:
                        continue
            body = "\n".join(body_chunks).strip()

            # Bilder extrahieren (falls vorhanden und noch Platz für weitere Bilder)
            if len(images) < MAX_IMAGES:
                for shp in slide.shapes:
                    if hasattr(shp, "image") and shp.image is not None:
                        try:
                            # Bild aus der Form extrahieren
                            img_data = shp.image.blob
                            img_ext = shp.image.ext
                            img_filename = f"slide_{i}_image_{len(images)+1}.{img_ext}"
                            img_path = tmpdir / img_filename

                            with open(img_path, "wb") as f:
                                f.write(img_data)

                            images.append(str(img_path))

                            if len(images) >= MAX_IMAGES:
                                break
                        except Exception:
                            # Bei Fehlern mit der nächsten Form fortfahren
                            continue

            # Notizen
            notes_txt = ""
            try:
                if (
                    getattr(slide, "has_notes_slide", False)
                    and slide.notes_slide is not None
                ):
                    ntf = getattr(slide.notes_slide, "notes_text_frame", None)
                    if ntf and getattr(ntf, "text", None):
                        notes_txt = ntf.text
            except Exception:
                # API-Varianz: notfalls ignorieren
                pass

            pieces.append(
                f"Slide {i}"
                + (f" — {title}" if title else "")
                + (f"\n{body}" if body else "")
                + (f"\nNotes: {notes_txt}" if notes_txt else "")
            )
            if len("\n\n".join(pieces)) > MAX_TEXT_CHARS * 1.5:
                break
    except Exception:
        # Bei Fehlern nur Text zurückgeben
        pass

    # Extrahierte Bilder skalieren
    resized_paths, resize_cleanup = resize_images_batch(
        images, max_dimension=MAX_IMAGE_DIMENSION
    )

    # Kombinierte Cleanup-Funktion
    def combined_cleanup():
        try:
            resize_cleanup()
        except Exception:
            pass
        try:
            tmpdir_obj.cleanup()
        except Exception:
            pass

    excerpt = _truncate("\n\n---\n\n".join(pieces))
    return excerpt, resized_paths, combined_cleanup


def ingest_file(path: Path) -> Tuple[str, List[str], callable]:
    ext = path.suffix.lower()
    if ext in {".jpg", ".jpeg", ".png"}:
        return ingest_image(path)
    if ext == ".docx":
        try:
            return ingest_docx(path)
        except Exception:
            return ingest_via_unstructured(path)
    if ext in {".doc", ".odt"}:
        return ingest_via_unstructured(path)
    if ext in {".pdf"}:
        return ingest_pdf(path)
    if ext in {".pptx"}:
        return ingest_pptx(path)
    if ext in {".ppt"}:
        return ingest_via_unstructured(path)
    if ext in {".txt"}:
        return ingest_txt(path)
    # Nicht unterstützte Dateiendungen werden als None zurückgegeben
    return None, [], lambda: None
